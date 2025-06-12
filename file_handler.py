import os
import tempfile
from typing import Optional
from werkzeug.utils import secure_filename
from pptx import Presentation
import fitz  # PyMuPDF
import docx
import re

ALLOWED_EXTENSIONS = {'pdf', 'pptx', 'docx', 'txt'}
MAX_FILE_SIZE = 16 * 1024 * 1024  # 16MB

class FileProcessingError(Exception):
    """Custom exception for file processing errors"""
    pass

def allowed_file(filename: str) -> bool:
    """Check if the file extension is allowed"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def check_file_size(file_size: int) -> bool:
    """Check if the file size is within limits"""
    return file_size <= MAX_FILE_SIZE

def get_file_extension(filename: str) -> str:
    """Get the file extension from filename"""
    return filename.rsplit('.', 1)[1].lower() if '.' in filename else ''

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file using PyMuPDF"""
    text = ""
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text() + "\n"
        doc.close()
        
        if not text.strip():
            raise FileProcessingError("No readable text found in PDF file")
    except Exception as e:
        raise FileProcessingError(f"Error processing PDF file: {str(e)}")
    return text

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint file"""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        if not text.strip():
            raise FileProcessingError("No readable text found in PowerPoint file")
    except Exception as e:
        raise FileProcessingError(f"Error processing PowerPoint file: {str(e)}")
    return text

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from Word document"""
    text = ""
    try:
        doc = docx.Document(file_path)
        for paragraph in doc.paragraphs:
            if paragraph.text:
                text += paragraph.text + "\n"
        if not text.strip():
            raise FileProcessingError("No readable text found in Word document")
    except Exception as e:
        raise FileProcessingError(f"Error processing Word document: {str(e)}")
    return text

def extract_text_from_txt(file_path: str) -> str:
    """Extract text from text file"""
    text = ""
    encodings = ['utf-8', 'latin-1', 'ascii']
    
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as file:
                text = file.read()
                if text.strip():
                    return text
        except UnicodeDecodeError:
            continue
        except Exception as e:
            raise FileProcessingError(f"Error processing text file: {str(e)}")
    
    if not text.strip():
        raise FileProcessingError("The text file is empty or contains invalid characters")
    return text

def clean_text(text: str) -> str:
    """Clean the extracted text"""
    # Remove extra whitespace and newlines
    text = re.sub(r'\s+', ' ', text)
    # Remove special characters but keep basic punctuation
    text = re.sub(r'[^\w\s.,!?-]', '', text)
    return text.strip()

def process_uploaded_file(file, upload_folder: str) -> str:
    """Process the uploaded file and extract text
    
    Args:
        file: The uploaded file object
        upload_folder: Path to the folder where files should be temporarily stored
        
    Returns:
        str: Extracted text from the file
        
    Raises:
        FileProcessingError: If there's an error processing the file
    """
    if not file or not file.filename:
        raise FileProcessingError("No file provided")
        
    filename = secure_filename(file.filename)
    if not allowed_file(filename):
        raise FileProcessingError(
            "File type not allowed. Please upload PDF, PPTX, DOCX, or TXT files only."
        )
    
    # Check file size
    file.seek(0, os.SEEK_END)
    size = file.tell()
    file.seek(0)
    
    if not check_file_size(size):
        raise FileProcessingError(
            f"File is too large. Maximum size is {MAX_FILE_SIZE / (1024 * 1024)}MB"
        )
    
    extension = get_file_extension(filename)
    
    # Create a temporary file to store the upload
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{extension}') as temp_file:
            file.save(temp_file.name)
            
            try:
                if extension == 'pdf':
                    text = extract_text_from_pdf(temp_file.name)
                elif extension == 'pptx':
                    text = extract_text_from_pptx(temp_file.name)
                elif extension == 'docx':
                    text = extract_text_from_docx(temp_file.name)
                elif extension == 'txt':
                    text = extract_text_from_txt(temp_file.name)
                else:
                    raise FileProcessingError("Unsupported file type")
                    
            finally:
                # Clean up the temporary file
                try:
                    os.unlink(temp_file.name)
                except Exception:
                    pass  # Ignore cleanup errors
                
    except Exception as e:
        raise FileProcessingError(f"Error saving or processing file: {str(e)}")
            
    if not text.strip():
        raise FileProcessingError("No text could be extracted from the file")
        
    # Clean and process the extracted text
    cleaned_text = clean_text(text)
    if not cleaned_text:
        raise FileProcessingError("No text content could be extracted from the file")
        
    return cleaned_text 